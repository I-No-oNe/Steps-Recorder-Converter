import os
import zipfile
from email import policy
from email.parser import BytesParser
from datetime import datetime
from PIL import Image
import io
import shutil

DOWNLOADS_DIR = os.path.join(os.path.expanduser("~"), "Downloads")
TEMP_EXTRACT_DIR = "temp_steps_extract"
TEMP_IMG_DIR = "temp_steps_images"


def clean_temp_dirs():
    for folder in [TEMP_EXTRACT_DIR, TEMP_IMG_DIR]:
        if os.path.exists(folder):
            shutil.rmtree(folder)
            print(f"Deleted temp folder: {folder}")


def find_latest_zip(folder):
    zips = [f for f in os.listdir(folder) if f.lower().endswith(".zip")]
    if not zips:
        raise FileNotFoundError("No zip files found in Downloads.")
    zips.sort(key=lambda x: os.path.getmtime(os.path.join(folder, x)), reverse=True)
    return os.path.join(folder, zips[0])


def extract_zip(zip_path, extract_to):
    with zipfile.ZipFile(zip_path, "r") as zip_ref:
        zip_ref.extractall(extract_to)


def find_mht_file(folder):
    for file in os.listdir(folder):
        if file.lower().endswith(".mht"):
            return os.path.join(folder, file)
    raise FileNotFoundError("No .mht file found after extraction.")


def extract_images_from_mht(mht_path):
    with open(mht_path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    images = []
    os.makedirs(TEMP_IMG_DIR, exist_ok=True)

    for part in msg.walk():
        content_type = part.get_content_type()
        if content_type.startswith("image/"):
            try:
                img_bytes = part.get_content()
                # Verify with Pillow
                img = Image.open(io.BytesIO(img_bytes))
                img.verify()

                ext = content_type.split("/")[1]
                img_name = f"image_{len(images)}.{ext}"
                img_path = os.path.join(TEMP_IMG_DIR, img_name)
                with open(img_path, "wb") as img_file:
                    img_file.write(img_bytes)
                images.append(img_path)
            except Exception as e:
                print(f"Skipping invalid image part: {e}")

    return images


def pixels_to_emu(px, dpi=96):
    # Convert pixels to English Metric Units (EMU)
    inches = px / dpi
    emu = int(inches * 914400)  # 1 inch = 914400 EMUs
    return emu


def create_presentation(image_paths):
    from pptx import Presentation

    prs = Presentation()

    for img_path in image_paths:
        # Open image to get dimensions
        with Image.open(img_path) as img:
            width_px, height_px = img.size
            dpi = img.info.get('dpi', (96, 96))[0]  # Use dpi from image if available, else 96
            width_emu = pixels_to_emu(width_px, dpi)
            height_emu = pixels_to_emu(height_px, dpi)

        # Set slide size to image size
        prs.slide_width = width_emu
        prs.slide_height = height_emu

        blank_slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Insert picture at (0,0) with original size
        slide.shapes.add_picture(img_path, 0, 0, width=width_emu, height=height_emu)

    output_file = f"StepsRecorder_{datetime.now().strftime('%Y-%m-%d_%H-%M-%S')}_sized.pptx"
    prs.save(output_file)
    print(f"Presentation saved: {output_file}")


def main():
    print("Cleaning temp directories...")
    clean_temp_dirs()

    print("Searching latest ZIP in Downloads...")
    zip_path = find_latest_zip(DOWNLOADS_DIR)
    print(f"Found ZIP: {zip_path}")

    os.makedirs(TEMP_EXTRACT_DIR, exist_ok=True)
    extract_zip(zip_path, TEMP_EXTRACT_DIR)

    mht_path = find_mht_file(TEMP_EXTRACT_DIR)
    print(f"Found MHT file: {mht_path}")

    print("Extracting images from MHT...")
    images = extract_images_from_mht(mht_path)
    print(f"Extracted {len(images)} images.")

    if images:
        print("Creating PowerPoint presentation...")
        create_presentation(images)
    else:
        print("No valid images extracted.")


if __name__ == "__main__":
    main()
